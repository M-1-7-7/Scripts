import os
import re
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from collections import defaultdict

# Function to find PowerPoint files
def find_powerpoint_files(directory):
    """
    This function finds all PowerPoint presentations (files with .ppt or .pptx extensions)
    within the specified directory and its subdirectories.

    Args:
        directory (str): Path to the directory to search.

    Returns:
        list: List containing the full paths to all PowerPoint presentations found.
    """
    powerpoint_files = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(('.ppt', '.pptx')):
                powerpoint_files.append(os.path.join(root, file))
    return powerpoint_files

# Function to load dirty words from file
def load_dirty_words(filename):
    """
    This function loads dirty words from a text file, removing any leading/trailing spaces
    and handling empty lines or an empty file gracefully.

    Args:
        filename (str): Path to the text file containing dirty words.

    Returns:
        list: List of cleaned dirty words (no leading/trailing spaces) or an empty list
              if the file is empty or contains only empty lines.
    """
    dirty_words = []
    try:
        with open(filename, 'r') as file:
            for line in file:
                # Remove leading/trailing spaces and add to list (if not empty)
                cleaned_word = line.strip()
                if cleaned_word:
                    dirty_words.append(cleaned_word.lower())
    except FileNotFoundError:
        print("Error: 'Dirty_word.txt' not found. Please ensure the file exists.")
    return dirty_words

# Load Dirty Words file
dirty_words = load_dirty_words('Dirty_word.txt')
if not dirty_words:
    print("Warning: 'Dirty_word.txt' is empty or contains only empty lines. No dirty words found.")

# Specify directory path (consider user input or relative paths)
computer_root = "C:\\"  # Replace with desired directory

powerpoint_files = find_powerpoint_files(computer_root)

# Dictionary to store findings for summary
findings_summary = defaultdict(lambda: defaultdict(int))

for powerpoint_file in powerpoint_files:
    try:
        prs = Presentation(powerpoint_file)

        for slide in prs.slides:
            # Process slide shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    shape_text = shape.text_frame.text.lower()

                    # Check for each dirty word using regular expressions for word boundaries
                    for word in dirty_words:
                        word_pattern = r'\b' + re.escape(word) + r'\b'
                        matches = re.findall(word_pattern, shape_text)
                        if matches:
                            findings_summary[(powerpoint_file, 'Slide Text')][word] += len(matches)

            # Process slide notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.lower()
                for word in dirty_words:
                    word_pattern = r'\b' + re.escape(word) + r'\b'
                    matches = re.findall(word_pattern, notes_text)
                    if matches:
                        findings_summary[(powerpoint_file, 'Presenter Notes')][word] += len(matches)

    except PermissionError:
        print(f"Insufficient system privileges to read contents of {powerpoint_file}")
    except PackageNotFoundError:
        print(f"Insufficient system privileges to read contents of {powerpoint_file}")
    except Exception as e:
        print(f"Error processing {powerpoint_file}: {e}")

# Print summary of findings
for (file, source), words in findings_summary.items():
    words_summary = ', '.join([f"{word}: {count}" for word, count in words.items()])
    print(f"File: {file} - {source} - Contains dirty words - {words_summary}")
