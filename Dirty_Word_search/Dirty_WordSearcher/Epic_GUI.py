import subprocess
import PyPDF2 
import os
import openpyxl
import re
import xlrd

import tkinter
from tkinter import *
from tkinter import ttk
from tkinter import messagebox

from spire.doc import *
from spire.doc.common import *
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from datetime import datetime
from collections import defaultdict

# Window Configuration
root = Tk()
root.geometry("800x600")
root.title("Dirty Word Searcher")
frm = ttk.Frame(root, padding=10)
frm.grid()
# user input variriables
DW_Var=tkinter.StringVar()
start_dir_Var=tkinter.StringVar()
out_dir_Var=tkinter.StringVar()
# buttons states
Excel_Button_Val = IntVar() 
PDF_Button_Val = IntVar()  
PowerPoint_Button_Val = IntVar() 
Doc_Button_Val = IntVar() 
TXT_Button_Val = IntVar() 
All_Other_Type_Button_Val = IntVar()
All_Document_Type_Button_Val = IntVar() 
# Other variables
script_directory = os.path.dirname(os.path.abspath(__file__)) 
# Global variables
global DW_Var_stripped
global out_dir_Var_stripped
global start_dir_Var_stripped

# Functions that will execute specificly to the checkboxes ticked by the user
def scan_all_docs():
    print("Scanning all docs")
    scan_excel_doc()
    scan_pdf_doc()
    scan_powerpoint_doc()
    scan_word_doc()
    scan_txt_doc()
    find_all_other_files()
# UPDATE THIS ONES
def scan_excel_doc():
    print("Scanning excel docs")
    positive_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Positive_Excel_Results.out"
    not_analysed_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Cannot_Anayse_Excel_List.out"  
    # Function to search for Excel files recursively in a directory
    def find_excel_files(directory):
        excel_files = []
        for root, dirs, files in os.walk(directory):
            for file in files:
                if file.lower().endswith(('.xls', '.xlsx')):
                    excel_files.append(os.path.join(root, file))
        return excel_files

    # Load dirty words from file
    Dirty_words = []
    with open(DW_Var.get(), 'r') as file:
        lines = file.readlines()
        for line in lines:
            Dirty_words.append(line.strip().lower())

    # Search for Excel files on the entire computer
    excel_files = find_excel_files(start_dir_Var.get())
    
    #create output directories for results
    if not os.path.exists(positive_output_file):
        with open(positive_output_file, "w") as file:
            file.write("THESE FILES HAVE ALL RETURNED A POSITIVE HIT\n------------\n------------\n")

    if not os.path.exists(not_analysed_output_file):
        with open(not_analysed_output_file, "w") as file:
            file.write("THIS FILE WHERE NOT ABLE TO BE ANALYSED BY THE SCRIPT\n------------\n------------\n")
    
    # Dictionary to store findings for summary
    findings_summary = defaultdict(lambda: defaultdict(int))

    def process_xlsx(file_path):
        """
        Process .xlsx files to find dirty words in all cells and comments.

        Args:
            file_path (str): Path to the .xlsx file to process.
        """
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
        except openpyxl.utils.exceptions.InvalidFileException:
            print(f"Error: Invalid .xlsx file {file_path}")
            return
        except Exception as e:
            print(f"Error opening .xlsx file {file_path}: {e}")
            return

        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            # Process cell text
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        cell_text = cell.value.lower()
                        for word in Dirty_words:
                            word_pattern = r'\b' + re.escape(word) + r'\b'
                            matches = re.findall(word_pattern, cell_text)
                            if matches:
                                findings_summary[(file_path, 'Cell Text')][word] += len(matches)
                    # Process comments
                    if cell.comment and cell.comment.text:
                        comment_text = cell.comment.text.lower()
                        for word in Dirty_words:
                            word_pattern = r'\b' + re.escape(word) + r'\b'
                            matches = re.findall(word_pattern, comment_text)
                            if matches:
                                findings_summary[(file_path, 'Comments')][word] += len(matches)

    def process_xls(file_path):
        """
        Process .xls files to find dirty words in all cells.

        Args:
            file_path (str): Path to the .xls file to process.
        """
        try:
            workbook = xlrd.open_workbook(file_path)
        except xlrd.biffh.XLRDError:
            print(f"Error: Invalid .xls file {file_path}")
            return
        except Exception as e:
            print(f"Error opening .xls file {file_path}: {e}")
            return

        for sheet in workbook.sheets():
            for row in range(sheet.nrows):
                for col in range(sheet.ncols):
                    cell_value = sheet.cell_value(row, col)
                    if isinstance(cell_value, str):
                        cell_text = cell_value.lower()
                        for word in Dirty_words:
                            word_pattern = r'\b' + re.escape(word) + r'\b'
                            matches = re.findall(word_pattern, cell_text)
                            if matches:
                                findings_summary[(file_path, 'Cell Text')][word] += len(matches)

    for excel_file in excel_files:
        try:
            if excel_file.lower().endswith('.xlsx'):
                process_xlsx(excel_file)
            elif excel_file.lower().endswith('.xls'):
                process_xls(excel_file)
        except PermissionError:
            print(f"Permission error processing {excel_file}")
        except Exception as e:
            print(f"Unexpected error processing {excel_file}: {e}")

    # Print summary of findings
    for (file, source), words in findings_summary.items():
        words_summary = ', '.join([f"{word}: {count}" for word, count in words.items()])
        pos_strings = f"File: {file} - {source} - Contains dirty words - {words_summary}"
        with open(positive_output_file, "a") as file:
            file.write(pos_strings + "\n")

def scan_pdf_doc():
    print("Scanning pdf docs")
    positive_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Positive_PDF_Results.out"
    no_txt_pdf_files = out_dir_Var.get() + "\\SCAN_RESULTS\\pdf_files_without_text.out"
    not_analysed_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Cannot_Anayse_PDF_List.out"  

    # Load dirty words from file
    Dirty_words = []
    with open(DW_Var.get(), 'r') as file:
        lines = file.readlines()
        for line in lines:
            Dirty_words.append(line.strip().lower())

    findings_summary = defaultdict(lambda: defaultdict(int))

    #create output directories for results
    if not os.path.exists(positive_output_file):
        with open(positive_output_file, "w") as file:
            file.write("THESE FILES HAVE ALL RETURNED A POSITIVE HIT\n------------\n------------\n")

    if not os.path.exists(no_txt_pdf_files):
        with open(no_txt_pdf_files, "w") as file:
            file.write("THESE FILES HAVE NO TXT AND CANNOT BE ANALYSED BY THE SCRIPT\n------------\n------------\n")

    if not os.path.exists(not_analysed_output_file):
        with open(not_analysed_output_file, "w") as file:
            file.write("THIS FILE WHERE NOT ABLE TO BE ANALYSED BY THE SCRIPT\n------------\n------------\n")

    # Function to search for PDF files recursively in a directory
    def find_pdf_files(directory):
        """
        This function finds all text files (files with .pdf extension)
        within the specified directory and its subdirectories.

        Args:
            directory (str): Path to the directory to search.

        Returns:
            list: List containing the full paths to all text files found.
        """
        pdf_files = []
        for root, dirs, files in os.walk(directory):
            for file in files:
                if file.lower().endswith('.pdf'):
                    pdf_files.append(os.path.join(root, file))
        return pdf_files

    def process_pdf_file(file_path):
        """
        Process text files to find dirty words.

        Args:
            file_path (str): Path to the text file to process.
        """
        # Perform search in PDF files
        try:
            # Open the PDF file
            with open(pdf_file, "rb"):
                reader = PyPDF2.PdfReader(pdf_file)
                    # Iterate over each page and extract text
                for page_num, page in enumerate(reader.pages, start=1):
                    text = page.extract_text()
                    # Split text into words and search for dirty words
                    words = text.split()
                    if len(text) > 0:
                        for word in Dirty_words:
                            for text_word in words:
                                word_pattern = r'\b' + re.escape(word) + r'\b'
                                matches = re.findall(word_pattern, text_word)
                                if matches:
                                    findings_summary[(file_path, 'Text')][word] += len(matches)
                    else:
                        print("This is a scanned Document")
                        string = f"File: {pdf_file}"
                        with open(no_txt_pdf_files, "a") as file:
                            file.write(string + "\n")
        except Exception as e:
            not_string = f"Error processing {pdf_file}: {e}"
            with open(not_analysed_output_file, "a") as file:
                file.write(not_string + "\n")

    # Search for PDF files on the entire computer
    pdf_files = find_pdf_files(start_dir_Var.get())

    # Loop through files in txt_files list
    for pdf_file in pdf_files:
        process_pdf_file(pdf_file)
    
    for (file, source), words in findings_summary.items():
        words_summary = ', '.join([f"{word}: {count}" for word, count in words.items()])
        pos_strings = f"File: {file} - {source} - Contains dirty words - {words_summary}"
        with open(positive_output_file, "a") as file:
            file.write(pos_strings + "\n")
    

def scan_powerpoint_doc():
    positive_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Positive_Powerpoint_Results.out"
    not_analysed_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Cannot_Anayse_Powerpoint_List.out"  
    def find_powerpoint_files(directory):
        powerpoint_files = []
        for root, dirs, files in os.walk(directory):
            for file in files:
                if file.lower().endswith(('.ppt', '.pptx')):
                    powerpoint_files.append(os.path.join(root, file))
        return powerpoint_files

    Dirty_words = []
    with open(DW_Var.get(), 'r') as file:
        lines = file.readlines()
        for line in lines:
            Dirty_words.append(line.strip().lower())

    powerpoint_files = find_powerpoint_files(start_dir_Var.get())

    #create output directories for results
    if not os.path.exists(positive_output_file):
        with open(positive_output_file, "w") as file:
            file.write("THESE FILES HAVE ALL RETURNED A POSITIVE HIT\n------------\n------------\n")

    if not os.path.exists(not_analysed_output_file):
        with open(not_analysed_output_file, "w") as file:
            file.write("THIS FILE WHERE NOT ABLE TO BE ANALYSED BY THE SCRIPT\n------------\n------------\n")

    # Dictionary to store findings for summary
    findings_summary = defaultdict(lambda: defaultdict(int))

    for powerpoint_file in powerpoint_files:
        try:
            prs = Presentation(powerpoint_file)

            for slide in prs.slides:
                # Process slide shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        shape_text = shape.text_frame.text.lower()

                        # Check for each dirty word using regular expressions for word boundaries
                        for word in Dirty_words:
                            word_pattern = r'\b' + re.escape(word) + r'\b'
                            matches = re.findall(word_pattern, shape_text)
                            if matches:
                                findings_summary[(powerpoint_file, 'Slide Text')][word] += len(matches)

                # Process slide notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.lower()
                    for word in Dirty_words:
                        word_pattern = r'\b' + re.escape(word) + r'\b'
                        matches = re.findall(word_pattern, notes_text)
                        if matches:
                            findings_summary[(powerpoint_file, 'Presenter Notes')][word] += len(matches)

        except PermissionError:
            print(f"Insufficient system privileges to read contents of {powerpoint_file}")
        except PackageNotFoundError:
            print(f"Insufficient system privileges to read contents of {powerpoint_file}")
        except Exception as e:
            print(f"Error processing {powerpoint_file}: {e}")

    # Print summary of findings
    for (file, source), words in findings_summary.items():
        words_summary = ', '.join([f"{word}: {count}" for word, count in words.items()])
        pos_strings = f"File: {file} - {source} - Contains dirty words - {words_summary}"
        with open(positive_output_file, "a") as file:
            file.write(pos_strings + "\n")

# Create error handling for scan_word_doc function
def scan_word_doc():
    #script_loc = script_directory + "\\word_searcher.py"
    print("Scanning word docs")
    #user input variables from GUI
    new_out_dir = out_dir_Var.get() + "\\SCAN_RESULTS\\Positive_Word_Results"
    not_analysed_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Cannot_Anayse_Word_List"  

    # Function to search for Word files recursively in a directory
    def find_word_files(directory):
        word_files = []
        for root, dirs, files in os.walk(directory):
            for file in files:
                if file.lower().endswith('.doc') or file.lower().endswith('.docx'):
                    word_files.append(os.path.join(root, file))
        return word_files

    #Create Dirty Word list from the dirty word input file from user
    Dirty_words = []
    with open(DW_Var.get(), 'r') as file:
        lines = file.readlines()
        for i in lines:
            Dirty_words.append(i.strip())
    
    # Create Positive folder if it doesn't exist
    if not os.path.exists(new_out_dir):
        os.makedirs(new_out_dir)

    # Get a list of Word files in the input directory
    word_files = find_word_files(start_dir_Var.get())

    # Process each Word file
    for eachfile in word_files:
        # Create an object of the Document class
        document = Document()
        # Load a Word document
        document.LoadFromFile(eachfile)
        # Find all instances of specific text
        for word in Dirty_words:
            textSelections = document.FindAllString(word, False, True)
            # Loop through all the instances
            for selection in textSelections:
                # Get the current instance as a single text range
                textRange = selection.GetAsOneRange()
                # Highlight the text range with a color
                textRange.CharacterFormat.HighlightColor = Color.get_Yellow()
                # Save the modified document
                output_file = os.path.join(new_out_dir, "POSITIVE - " + os.path.basename(eachfile))
                document.SaveToFile(output_file)
        # Close the document
        document.Close()
    print("finished scan")

def scan_txt_doc():
    print("scanning txt docs")
    positive_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Positive_TXT_Results.out"
    not_analysed_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Cannot_Anayse_TXT_List.out"  
    
    # load dirty words from file into Dirty_words list
    Dirty_words = []
    with open(DW_Var.get(), 'r') as file:
        lines = file.readlines()
        for line in lines:
            Dirty_words.append(line.strip().lower())

    findings_summary = defaultdict(lambda: defaultdict(int))

    #create output directories/files for results
    if not os.path.exists(positive_output_file):
        with open(positive_output_file, "w") as file:
            file.write("THESE FILES HAVE ALL RETURNED A POSITIVE HIT\n------------\n------------\n")

    if not os.path.exists(not_analysed_output_file):
        with open(not_analysed_output_file, "w") as file:
            file.write("THIS FILE WHERE NOT ABLE TO BE ANALYSED BY THE SCRIPT\n------------\n------------\n")

    # This funtion will be called and pass through the files nested within the target directory for files with txt extension
    def find_text_files(directory):
        """
        This function finds all text files (files with .txt extension)
        within the specified directory and its subdirectories.

        Args:
            directory (str): Path to the directory to search.

        Returns:
            list: List containing the full paths to all text files found.
        """
        text_files = []
        for root, dirs, files in os.walk(directory):
            for file in files:
                if file.lower().endswith('.txt'):
                    text_files.append(os.path.join(root, file))
        return text_files

    # Read the contents of txt files discovered
    def read_file(file_path):
        """
        Attempt to read a file with multiple encodings.
        
        Args:
            file_path (str): Path to the file to read.

        Returns:
            str: Decoded file content, or None if decoding fails.
        """
        encodings = ['utf-8', 'latin-1', 'ISO-8859-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read().lower()
            except UnicodeDecodeError:
                continue
        return None

    # Identify matches to words in Dirty_word list
    def process_text_file(file_path):
        """
        Process text files to find dirty words.

        Args:
            file_path (str): Path to the text file to process.
        """
        try:
            file_text = read_file(file_path)
            if file_text is None:
                neg_string = f"Error: Could not decode {file_path}"
                with open(not_analysed_output_file, "a") as file:
                    file.write(neg_string + "\n")
                return
            for word in Dirty_words:
                word_pattern = r'\b' + re.escape(word) + r'\b'
                matches = re.findall(word_pattern, file_text)
                if matches:
                    findings_summary[(file_path, 'Text')][word] += len(matches)
        except FileNotFoundError:
            neg_string = f"Error: File not found {file_path}"
            with open(not_analysed_output_file, "a") as file:
                file.write(neg_string + "\n")
        except PermissionError:
            neg_string = f"Permission error processing {file_path}"
            with open(not_analysed_output_file, "a") as file:
                file.write(neg_string + "\n")
        except Exception as e:
            neg_string = f"Unexpected error processing {file_path}: {e}"
            with open(not_analysed_output_file, "a") as file:
                file.write(neg_string + "\n")

    # start searching for txt files
    text_files = find_text_files(start_dir_Var.get())

    # Loop through files in txt_files list
    for text_file in text_files:
        process_text_file(text_file)
    
    for (file, source), words in findings_summary.items():
        words_summary = ', '.join([f"{word}: {count}" for word, count in words.items()])
        pos_strings = f"File: {file} - {source} - Contains dirty words - {words_summary}"
        with open(positive_output_file, "a") as file:
                        file.write(pos_strings + "\n")
def find_all_other_files():
    positive_output_file = out_dir_Var.get() + "\\SCAN_RESULTS\\Other_Files"
    start_dir = start_dir_Var.get()
    print("finding all other files")
    # Lists for all the file types
    list_of_lists = {
    "jpeg": [], 
    "png": [],
    "bmp": [],
    "tiff": [],
    "mp3": [],
    "mp4": [],
    "avi": [],
    "mov": [],
    "msg": []
    }
     # Create directory for the otherfiles output
    if not os.path.exists(positive_output_file):
        os.mkdir(positive_output_file)
    
    interesting_extensions = [".jpeg", ".png", ".bmp", ".tiff", ".mp3", ".mp4", ".avi", ".mov", ".msg"]
    
    for i in interesting_extensions:
        targer_array = i.replace(".", "")
        for root, dirs, files in os.walk(start_dir):
            for file in files:
                if file.lower().endswith(i):
                    file_found = os.path.join(root, file)
                    list_of_lists[targer_array].append(file_found)

    if len(list_of_lists["jpeg"]) > 1:
        file_list = positive_output_file + "\\jpegs_found.out"
        if not os.path.exists(file_list):
            with open(file_list, "w") as file:
                file.write("THESE FILES HAVE BEEN FOUND\n------------\n------------\n")

        for i in list_of_lists["jpeg"]:
            with open(file_list, "a") as file:
                    file.write(i + "\n")

    if len(list_of_lists["png"]) > 1:
        file_list = positive_output_file + "\\png_found.out"
        if not os.path.exists(file_list):
            with open(file_list, "w") as file:
                file.write("THESE FILES HAVE BEEN FOUND\n------------\n------------\n")

        for i in list_of_lists["png"]:
            with open(file_list, "a") as file:
                    file.write(i + "\n")

    if len(list_of_lists["bmp"]) > 1:
        file_list = positive_output_file + "\\bmp_found.out"
        if not os.path.exists(file_list):
            with open(file_list, "w") as file:
                file.write("THESE FILES HAVE BEEN FOUND\n------------\n------------\n")

        for i in list_of_lists["bmp"]:
            with open(file_list, "a") as file:
                    file.write(i + "\n")
    
    if len(list_of_lists["tiff"]) > 1:
        file_list = positive_output_file + "\\tiff_found.out"
        if not os.path.exists(file_list):
            with open(file_list, "w") as file:
                file.write("THESE FILES HAVE BEEN FOUND\n------------\n------------\n")

        for i in list_of_lists["tiff"]:
            with open(file_list, "a") as file:
                    file.write(i + "\n")

    if len(list_of_lists["mp3"]) > 1:
        file_list = positive_output_file + "\\mp3_found.out"
        if not os.path.exists(file_list):
            with open(file_list, "w") as file:
                file.write("THESE FILES HAVE BEEN FOUND\n------------\n------------\n")

        for i in list_of_lists["mp3"]:
            with open(file_list, "a") as file:
                    file.write(i + "\n")
    
    if len(list_of_lists["mp4"]) > 1:
        file_list = positive_output_file + "\\mp4_found.out"
        if not os.path.exists(file_list):
            with open(file_list, "w") as file:
                file.write("THESE FILES HAVE BEEN FOUND\n------------\n------------\n")

        for i in list_of_lists["mp4"]:
            with open(file_list, "a") as file:
                    file.write(i + "\n")
    
    if len(list_of_lists["avi"]) > 1:
        file_list = positive_output_file + "\\avi_found.out"
        if not os.path.exists(file_list):
            with open(file_list, "w") as file:
                file.write("THESE FILES HAVE BEEN FOUND\n------------\n------------\n")

        for i in list_of_lists["avi"]:
            with open(file_list, "a") as file:
                    file.write(i + "\n")

    if len(list_of_lists["mov"]) > 1:
        file_list = positive_output_file + "\\mov_found.out"
        if not os.path.exists(file_list):
            with open(file_list, "w") as file:
                file.write("THESE FILES HAVE BEEN FOUND\n------------\n------------\n")

        for i in list_of_lists["mov"]:
            with open(file_list, "a") as file:
                    file.write(i + "\n")
    
    if len(list_of_lists["msg"]) > 1:
        file_list = positive_output_file + "\\msg_found.out"
        if not os.path.exists(file_list):
            with open(file_list, "w") as file:
                file.write("THESE FILES HAVE BEEN FOUND\n------------\n------------\n")

        for i in list_of_lists["msg"]:
            with open(file_list, "a") as file:
                    file.write(i + "\n")

    print("On to the next")

# Functions from GUI Buttons
def search_files():
    output_dir = out_dir_Var.get() + "\\SCAN_RESULTS"
    if not os.path.exists(output_dir):
        os.mkdir(output_dir)
    # Ensure the files and Dirs entered by user are going to work
    if All_Document_Type_Button_Val.get() == 1:
        print("Scanning all docs")
        scan_excel_doc()
        scan_pdf_doc()
        scan_powerpoint_doc()
        scan_word_doc()
        scan_txt_doc()

    elif All_Document_Type_Button_Val.get() == 0:
        if Excel_Button_Val.get() == 1:
            scan_excel_doc()

        if PDF_Button_Val.get() == 1:
            scan_pdf_doc()

        if PowerPoint_Button_Val.get() == 1:
            scan_powerpoint_doc()

        if Doc_Button_Val.get() == 1:
            scan_word_doc()

        if TXT_Button_Val.get() == 1:
            scan_txt_doc()

        if All_Other_Type_Button_Val.get() == 1:
            find_all_other_files()

    tkinter.messagebox.showinfo(title="Status", message="Files have been analysed")

# Validate user selections and inputs 
def check_checkboxes():
    if All_Document_Type_Button_Val.get() == 0 and PDF_Button_Val.get() == 0 and PowerPoint_Button_Val.get() == 0 and Doc_Button_Val.get() == 0 and TXT_Button_Val.get() == 0 and Excel_Button_Val.get() == 0 and All_Other_Type_Button_Val.get() == 0:
        messagebox.showerror("Checkbox Status", "please check at least one checkbox")
        main_screen()
    else:
        new_out_dir = out_dir_Var.get() + "\\SCAN_RESULTS\\"
        if not os.path.exists(new_out_dir):
            os.makedirs(new_out_dir)
        search_files()

def check_user_input():
    user_input_bad = ""
    global start_dir_Var_stripped
    global out_dir_Var_stripped
    global DW_Var_stripped
    
    # Stip user inputs
    start_dir_Var_stripped = start_dir_Var.get()
    start_dir_Var_stripped = start_dir_Var_stripped.strip()
    out_dir_Var_stripped = out_dir_Var.get()
    out_dir_Var_stripped = out_dir_Var_stripped.strip()
    DW_Var_stripped = DW_Var.get()
    DW_Var_stripped = DW_Var_stripped.strip()

    # check if directory exists
    if not os.path.exists(start_dir_Var_stripped):
        user_input_bad += "!!!Please enter a valid STARTING DIRECTORY!!!\n"
    else:
        print("starting directory is existing :)\n")

    if not os.path.exists(out_dir_Var_stripped):
        user_input_bad += "!!!Please enter a valid OUTPUT DIRECTORY!!!\n"
    else:
       print("Output directory is existing :)\n")

    # Check if the file is writable using os.access()
    if DW_Var_stripped.endswith('.txt'):
        if os.access(DW_Var_stripped, os.W_OK):
            print(f"File '{DW_Var_stripped}' is writable.") 
        elif FileNotFoundError:
            user_input_bad += "!!!Dirty Word File NOT FOUND!!!\n"
        elif PermissionError:
            user_input_bad += "!!!Please enter a valid FILE WHERE PROGRAM HAS PERMISIONS!!!\n"
    else:
        user_input_bad += "!!!Please enter a valid TXT FILE!!!\n"

    if len(user_input_bad) > 0:
        messagebox.showerror("Your Input Status", user_input_bad)
    else:
        check_checkboxes()

#reset values on the window
def reset_values():
    print("Resetting")
    DW_Var.set("")
    start_dir_Var.set("")
    out_dir_Var.set("")
    Excel_Button_Val.set(0) 
    PDF_Button_Val.set(0) 
    PowerPoint_Button_Val.set(0) 
    Doc_Button_Val.set(0) 
    TXT_Button_Val.set(0) 
    All_Document_Type_Button_Val.set(0) 

# Main Screen
def main_screen():
    title = ttk.Label(frm, text="WELCOME TO DIRTY WORD SEARCHER!!!! :)")
    # User input fields
    DW_label = ttk.Label(frm, text = "Dirty Word Text file: ")
    DW_entry = ttk.Entry(frm, textvariable = DW_Var)
    start_dir_label = ttk.Label(frm, text = "Starting Directory for your search : ")
    start_dir_entry = ttk.Entry(frm, textvariable = start_dir_Var)
    out_dir_label = ttk.Label(frm, text = "Output Directory for your results : ")
    out_dir_entry = ttk.Entry(frm, textvariable = out_dir_Var)
    # Radio Buttons for file types

    # Function Buttons
    #show_doc_results = ttk.Button(frm, text="Show Doc Results", command=)
    #show_pptx_results = ttk.Button(frm, text="Show pptx Results", command=)
    #show_pdf_results = ttk.Button(frm, text="Show pdf Results", command=)
    #show_xls_results = ttk.Button(frm, text="Show xls Results", command=)
    search_btn = ttk.Button(frm, text="Search", command=check_user_input)
    reset_btn = ttk.Button(frm, text="Reset", command=reset_values)
    quit_btn = ttk.Button(frm, text="Quit", command=root.destroy)

    All_Excel_Button = Checkbutton(root, text = "All EXCEL Docs", 
                        variable = Excel_Button_Val, 
                        onvalue = 1, 
                        offvalue = 0, 
                        height = 2, 
                        width = 20) 
                                              
    All_PDF_Button = Checkbutton(root, text = "All PDF Docs", 
                        variable = PDF_Button_Val, 
                        onvalue = 1, 
                        offvalue = 0, 
                        height = 2, 
                        width = 20) 
                                                     
    All_PowerPoint_Button = Checkbutton(root, text = "All PowerPoint Docs", 
                        variable = PowerPoint_Button_Val, 
                        onvalue = 1, 
                        offvalue = 0, 
                        height = 2, 
                        width = 20) 
                                               
    All_Word_Button = Checkbutton(root, text = "All Word Docs", 
                        variable = Doc_Button_Val, 
                        onvalue = 1, 
                        offvalue = 0, 
                        height = 2, 
                        width = 20) 
                                              
    All_TXT_Button = Checkbutton(root, text = "All TXT Docs", 
                        variable = TXT_Button_Val, 
                        onvalue = 1, 
                        offvalue = 0, 
                        height = 2, 
                        width = 20) 

    All_Documet_Type_Button = Checkbutton(root, text = "All Document Types", 
                        variable = All_Document_Type_Button_Val, 
                        onvalue = 1, 
                        offvalue = 0, 
                        height = 2, 
                        width = 20) 
    All_Other_Type_Button = Checkbutton(root, text = "All Other Types", 
                        variable = All_Other_Type_Button_Val, 
                        onvalue = 1, 
                        offvalue = 0, 
                        height = 2, 
                        width = 20) 
    
   
    # Placement of elements on frame
    title.grid(column=0, row=0)

    DW_label.grid(column=0, row=3, sticky=E)
    DW_entry.grid(column=1, row=3, columnspan=4)
    start_dir_label.grid(column=0, row=4, sticky=E)
    start_dir_entry.grid(column=1, row=4, columnspan=4)
    out_dir_label.grid(column=0, row=5, sticky=E)
    out_dir_entry.grid(column=1, row=5, columnspan=4)
    
    All_Excel_Button.grid(row=2, sticky=W)
    All_PDF_Button.grid(row=3, sticky=W)
    All_PowerPoint_Button.grid(row=4, sticky=W)
    All_Word_Button.grid(row=5, sticky=W)
    All_TXT_Button.grid(row=6, sticky=W)
    All_Other_Type_Button.grid(row=7, sticky=W)
    All_Documet_Type_Button.grid(row=8, sticky=W)

    search_btn.grid(column=8, row=10)
    reset_btn.grid(column=9, row=10)
    quit_btn.grid(column=10, row=10)

    root.mainloop()

if __name__ == "__main__":
    main_screen()
